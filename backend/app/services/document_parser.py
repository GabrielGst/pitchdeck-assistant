"""
DocumentParser — extracts text from PDF and PPTX files.

Primary seam: parse(file_bytes, mime_type) → ExtractedDeck
Each slide/page becomes a Chunk with its index and raw text.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field


@dataclass
class Chunk:
    index: int       # slide number (PPTX) or page number (PDF), 0-based
    text: str


@dataclass
class ExtractedDeck:
    chunks: list[Chunk] = field(default_factory=list)
    company_name: str = ""

    @property
    def full_text(self) -> str:
        return "\n\n".join(c.text for c in self.chunks if c.text.strip())


SUPPORTED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def parse(file_bytes: bytes, mime_type: str, filename: str = "") -> ExtractedDeck:
    """Extract text chunks from a PDF or PPTX file."""
    if mime_type not in SUPPORTED_MIME_TYPES:
        raise ValueError(f"Unsupported file type: {mime_type}")

    if mime_type == "application/pdf":
        return _parse_pdf(file_bytes, filename)
    return _parse_pptx(file_bytes, filename)


def _parse_pdf(file_bytes: bytes, filename: str) -> ExtractedDeck:
    import pdfplumber

    chunks: list[Chunk] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            chunks.append(Chunk(index=i, text=text.strip()))

    return ExtractedDeck(chunks=chunks, company_name=_infer_company_name(chunks, filename))


def _parse_pptx(file_bytes: bytes, filename: str) -> ExtractedDeck:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: list[Chunk] = []

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text)
                    if line.strip():
                        texts.append(line.strip())
        chunks.append(Chunk(index=i, text="\n".join(texts)))

    return ExtractedDeck(chunks=chunks, company_name=_infer_company_name(chunks, filename))


def _infer_company_name(chunks: list[Chunk], filename: str) -> str:
    """Best-effort: first non-empty text block from slide 0, fallback to filename stem."""
    if chunks:
        first_lines = [ln for ln in chunks[0].text.splitlines() if ln.strip()]
        if first_lines:
            return first_lines[0][:200]
    # Strip extension from filename
    stem = filename.rsplit(".", 1)[0] if "." in filename else filename
    return stem or "Unknown"
