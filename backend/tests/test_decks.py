"""
Tests for POST /decks (upload) and GET /deals/kanban — primary API seam.
External dependencies (Celery, DB) are mocked.
"""

import io
import uuid
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from pptx.util import Inches

from app.main import app
from app.models.base import Deal, DealStage, Deck, DeckStatus, Role, Tenant, User

MOCK_TENANT = Tenant(id=uuid.uuid4(), name="Test VC", slug="test-vc")
MOCK_USER = User(
    id=uuid.uuid4(),
    tenant_id=MOCK_TENANT.id,
    clerk_user_id="user_test123",
    email="analyst@testvc.com",
    role=Role.analyst,
)


def _make_pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    txBox.text_frame.text = "Test Company"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def authenticated_client():
    """Client with auth mocked to return MOCK_USER."""
    async def _get():
        async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as ac:
            yield ac

    return _get


@pytest.fixture
async def client():
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as ac:
        yield ac


async def test_upload_rejects_unauthenticated(client: AsyncClient) -> None:
    pptx_bytes = _make_pptx_bytes()
    resp = await client.post(
        "/decks",
        files={"file": ("deck.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert resp.status_code == 401


async def test_upload_rejects_unsupported_type(client: AsyncClient) -> None:
    mock_result = MagicMock()
    mock_result.scalar_one_or_none.return_value = MOCK_USER
    mock_session = AsyncMock()
    mock_session.execute = AsyncMock(return_value=mock_result)
    mock_session.get = AsyncMock(return_value=MOCK_TENANT)
    mock_session.__aenter__ = AsyncMock(return_value=mock_session)
    mock_session.__aexit__ = AsyncMock(return_value=None)

    with (
        patch("app.api.deps.verify_clerk_token", return_value="user_test123"),
        patch("app.core.database.AsyncSessionLocal", return_value=mock_session),
    ):
        resp = await client.post(
            "/decks",
            files={"file": ("report.docx", b"content", "application/msword")},
            headers={"Authorization": "Bearer token"},
        )
    assert resp.status_code == 422
    assert "Unsupported" in resp.json()["detail"]


async def test_kanban_returns_all_stages(client: AsyncClient) -> None:
    """GET /deals/kanban returns a list with all 6 pipeline stages."""
    mock_result = MagicMock()
    mock_result.scalar_one_or_none.return_value = MOCK_USER
    mock_deals_result = MagicMock()
    mock_deals_result.scalars.return_value.all.return_value = []
    mock_session = AsyncMock()
    mock_session.execute = AsyncMock(side_effect=[mock_result, mock_deals_result])
    mock_session.get = AsyncMock(return_value=MOCK_TENANT)
    mock_session.__aenter__ = AsyncMock(return_value=mock_session)
    mock_session.__aexit__ = AsyncMock(return_value=None)

    with (
        patch("app.api.deps.verify_clerk_token", return_value="user_test123"),
        patch("app.core.database.AsyncSessionLocal", return_value=mock_session),
    ):
        resp = await client.get("/deals/kanban", headers={"Authorization": "Bearer token"})

    assert resp.status_code == 200
    stages = {col["stage"] for col in resp.json()}
    assert stages == {"inbox", "screening", "due_diligence", "partner_review", "invested", "passed"}
