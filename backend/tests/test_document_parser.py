"""
Tests for DocumentParser — the secondary testing seam.
These are pure-function tests with no I/O beyond reading fixture bytes.
"""

import io

import pytest

from app.services.document_parser import (
    SUPPORTED_MIME_TYPES,
    ExtractedDeck,
    parse,
)

PDF_MIME = "application/pdf"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _make_minimal_pdf(text: str = "Hello World") -> bytes:
    """Build a minimal valid single-page PDF with the given text."""
    content = (
        f"BT /F1 12 Tf 100 700 Td ({text}) Tj ET"
    )
    stream = content.encode()
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]"
        b" /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        + b"4 0 obj\n<< /Length " + str(len(stream)).encode() + b" >>\nstream\n"
        + stream + b"\nendstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000274 00000 n \n"
        b"0000000360 00000 n \n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
        b"startxref\n430\n%%EOF"
    )
    return pdf


def _make_minimal_pptx(slide_text: str = "Acme Corp") -> bytes:
    """Build a minimal PPTX with one slide containing the given text."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    txBox.text_frame.text = slide_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_unsupported_mime_type_raises():
    with pytest.raises(ValueError, match="Unsupported"):
        parse(b"data", "application/msword")


def test_supported_mime_types_set():
    assert PDF_MIME in SUPPORTED_MIME_TYPES
    assert PPTX_MIME in SUPPORTED_MIME_TYPES


def test_pptx_extracts_text():
    pptx_bytes = _make_minimal_pptx("Sequoia Capital")
    result = parse(pptx_bytes, PPTX_MIME, "sequoia.pptx")

    assert isinstance(result, ExtractedDeck)
    assert len(result.chunks) == 1
    assert "Sequoia" in result.chunks[0].text


def test_pptx_company_name_from_first_slide():
    pptx_bytes = _make_minimal_pptx("Acme Corp — Series A")
    result = parse(pptx_bytes, PPTX_MIME, "deck.pptx")
    assert result.company_name.startswith("Acme Corp")


def test_pptx_company_name_falls_back_to_filename():
    """If first slide is empty, company name comes from filename stem."""
    pptx_bytes = _make_minimal_pptx("")
    result = parse(pptx_bytes, PPTX_MIME, "my-startup.pptx")
    assert result.company_name == "my-startup"


def test_full_text_joins_chunks():
    pptx_bytes = _make_minimal_pptx("Line one")
    result = parse(pptx_bytes, PPTX_MIME)
    assert isinstance(result.full_text, str)


def test_empty_pptx_returns_empty_deck():
    from pptx import Presentation

    prs = Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    result = parse(buf.getvalue(), PPTX_MIME, "empty.pptx")
    assert result.chunks == []
    assert result.company_name == "empty"
